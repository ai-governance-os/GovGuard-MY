"""Real Office tool handlers: docx, pptx, xlsx — with safety guard."""
from __future__ import annotations

from ..models import CandidateAction
from ..util.path_guard import resolve_safe
from ._safety import safe_target_check


class DocxTool:
    name = "docx"

    def __init__(self, workspace_roots: list[str]) -> None:
        self.roots = list(workspace_roots)

    def __call__(self, action: CandidateAction) -> dict:
        ok, reason = safe_target_check(action.target, self.roots)
        if not ok:
            return _denied(reason)
        try:
            from docx import Document  # type: ignore
        except ImportError:
            return _failed("python-docx not installed")
        path = resolve_safe(action.target)
        path.parent.mkdir(parents=True, exist_ok=True)
        doc = Document()
        meta = action.metadata or {}
        title = meta.get("title") or action.purpose
        if title:
            doc.add_heading(title, level=0)
        for heading in meta.get("headings", []):
            doc.add_heading(str(heading), level=1)
        body = meta.get("body") or meta.get("content") or ""
        for paragraph in str(body).split("\n\n"):
            if paragraph.strip():
                doc.add_paragraph(paragraph.strip())
        doc.save(str(path))
        return {"status": "success", "summary": f"docx_written:{path.name}", "affected": [str(path)]}


class PptxTool:
    name = "pptx"

    def __init__(self, workspace_roots: list[str]) -> None:
        self.roots = list(workspace_roots)

    def __call__(self, action: CandidateAction) -> dict:
        ok, reason = safe_target_check(action.target, self.roots)
        if not ok:
            return _denied(reason)
        try:
            from pptx import Presentation  # type: ignore
        except ImportError:
            return _failed("python-pptx not installed")
        path = resolve_safe(action.target)
        path.parent.mkdir(parents=True, exist_ok=True)
        prs = Presentation()
        meta = action.metadata or {}
        title_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_layout)
        slide.shapes.title.text = meta.get("title") or action.purpose or "Presentation"
        if slide.placeholders[1] is not None:
            slide.placeholders[1].text = meta.get("subtitle", "")
        bullet_layout = prs.slide_layouts[1]
        for s in meta.get("slides", []):
            slide = prs.slides.add_slide(bullet_layout)
            slide.shapes.title.text = str(s.get("title", ""))
            tf = slide.placeholders[1].text_frame
            tf.text = ""
            for i, b in enumerate(s.get("bullets", [])):
                if i == 0:
                    tf.text = str(b)
                else:
                    tf.add_paragraph().text = str(b)
        prs.save(str(path))
        return {"status": "success", "summary": f"pptx_written:{path.name}", "affected": [str(path)]}


class XlsxTool:
    name = "xlsx"

    def __init__(self, workspace_roots: list[str]) -> None:
        self.roots = list(workspace_roots)

    def __call__(self, action: CandidateAction) -> dict:
        ok, reason = safe_target_check(action.target, self.roots)
        if not ok:
            return _denied(reason)
        try:
            from openpyxl import Workbook  # type: ignore
        except ImportError:
            return _failed("openpyxl not installed")
        path = resolve_safe(action.target)
        path.parent.mkdir(parents=True, exist_ok=True)
        wb = Workbook()
        wb.remove(wb.active)
        meta = action.metadata or {}
        sheets = meta.get("sheets") or {"Sheet1": meta.get("rows", [])}
        for sheet_name, rows in sheets.items():
            ws = wb.create_sheet(title=str(sheet_name)[:31])
            for row in rows:
                ws.append(list(row))
        if not wb.sheetnames:
            wb.create_sheet(title="Sheet1")
        wb.save(str(path))
        return {"status": "success", "summary": f"xlsx_written:{path.name}", "affected": [str(path)]}


def _denied(reason: str) -> dict:
    return {"status": "denied", "summary": reason, "affected": [], "error": reason}


def _failed(reason: str) -> dict:
    return {"status": "failed", "summary": reason, "affected": [], "error": reason}
